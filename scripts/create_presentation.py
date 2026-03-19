from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()

def add_slide(title_text, content_text="", placeholder_text=None):
    layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    title.text = title_text
    
    # Adjust title font
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 114, 178) # Okabe-Ito Blue
            
    content = slide.placeholders[1]
    content.text = content_text
    
    # Large, easily readable font for sparse keywords
    for paragraph in content.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(28)
            
    if placeholder_text:
        # Add a colored box to indicate exactly where the explicit dashboard screenshot should go
        txBox = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(6), Inches(3))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"[ PASTE DASHBOARD JPG HERE: \n{placeholder_text} ]"
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(26)
        p.font.color.rgb = RGBColor(213, 94, 0) # Okabe-Ito Vermillion
        
    return slide

# Slide 1: Title
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Using Optic Fibers as Weather Sensors"
subtitle.text = "Group Members: [Your Names Here]\nData Mining & Analytics"

# Update Title Font
for p in title.text_frame.paragraphs:
    for run in p.runs:
        run.font.name = "Arial"
        run.font.color.rgb = RGBColor(0, 114, 178)

# Slide 2
add_slide("Project Description", 
"• Location: Erfurt to Sundhausen fiber link.\n"
"• Concept: Laser light has a specific shape ('Polarization').\n"
"• Problem: Outside weather physically twists the fiber glass.\n"
"• Goal: Can we detect the weather simply by looking at the light?")

# Slide 3
add_slide("Data Sources", 
"• Optical Data: PAX1000 Polarimeter.\n"
"  - Extremely fast (reads 4 times a second).\n"
"• Weather Data: Open-Meteo API.\n"
"  - Slower logs (Air Pressure, Temperature, Humidity).")

# Slide 4
add_slide("Cleaning the Data", 
"• Fix broken sensors and missing timestamps.\n"
"• Remove absolute machine errors (-99.990).\n"
"• Match the fast light data perfectly with the slow weather data.")

# Slide 5
add_slide("Active Data Metrics", 
"• Total Raw Data: ~2.37 Million points.\n"
"• Bad Data Removed: ~1.12 Million anomalies.\n"
"• Clean Data Used: ~1.25 Million valid points.",
"Tab 4 'Data Cleaning Summary'")

# Slide 6
add_slide("Tracking Light Over Time", 
"• Blue Line: Azimuth (The 'tilt' of the light).\n"
"• Orange Line: Ellipticity (How 'circular' the light is).\n"
"• Finding: The light constantly swings back and forth over days.",
"Tab 1 'Azimuth & Ellipticity Combined'")

# Slide 7
add_slide("Watching the Shape Change", 
"• The light forms an ellipse (an oval shape).\n"
"• When the weather changes, the oval stretches and rotates.\n"
"• Our dashboard animates this exact movement visually.",
"Tab 3 'Live Ellipse Preview'")

# Slide 8
add_slide("The Effect of Surface Pressure",
"• Question: Does heavy air pressure crush the glass?\n"
"• Graph: Air Pressure (X-axis) vs. Light Angles (Y-axis).\n"
"• Red Trendline: A regression line calculating the pattern in one stroke.\n"
"• Finding: Yes! High pressure strictly locks the light into specific angles.",
"Tab 2 'Azimuth vs Pressure' & 'Ellipticity vs Pressure'")

# Slide 9
add_slide("Discovering Daily Rhythms",
"• Question: Does the sun affect the fiber every single day?\n"
"• Method: Average the last 3 days of data by hour (0–23h).\n"
"• Finding: A clear daily wave appears—the sun creates a predictable heartbeat!",
"Tab 1 '24-Hour Daily Rhythms' chart (top right)")

# Slide 10
add_slide("Conclusion", 
"• Primary Finding: Aerial fibers act completely like giant weather sensors.\n"
"• Physical Interaction: Heavy air pressure visibly warps the geometric data inside the glass.\n"
"• Final Takeaway: We can theoretically turn existing telecommunication cables into tracking networks.")

os.makedirs("presentation", exist_ok=True)
prs.save("presentation/Final_Project_Presentation_V7.pptx")
print("Presentation compiled successfully.")
